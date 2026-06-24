"""PPT + 交互式 HTML 导出模块。

生成:
  - PPTX: 封面/答案/指标/chunks/可视化截图 (线性浏览，适合演示)
  - HTML: 自包含交互式页面 (含 zoom/pan 图表，浏览器打开探索)

Usage:
    from frontend.utils.ppt_exporter import build_pptx_bytes, build_interactive_html
"""

from __future__ import annotations

import base64
import io
import re
import sys
import tempfile
from datetime import datetime
from pathlib import Path

# Streamlit 直接执行页面文件时不设置包上下文，需自行确保项目根在 sys.path
_PROJECT_ROOT = str(Path(__file__).parent.parent.parent)
if _PROJECT_ROOT not in sys.path:
    sys.path.insert(0, _PROJECT_ROOT)

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from frontend.styles import DEPTH_COLORS

# ── Playwright 可选依赖 ──
try:
    from playwright.sync_api import sync_playwright
    HAS_PLAYWRIGHT = True
except ImportError:
    HAS_PLAYWRIGHT = False

# ── 常量 ──
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FONT_CN = "Microsoft YaHei"
FONT_EN = "Calibri"
ACCENT = DEPTH_COLORS[0]  # Deep blue #2563EB
ALGO_LABELS = {
    "llm-only": "LLM Only", "naive-rag": "模块化 RAG",
    "rag-with-judge": "递归检索 RAG", "rag-loop": "规划-执行-反馈闭环 RAG",
}
SOURCE_LABELS = {"milvus": "Milvus 向量数据库", "web": "网络搜索 (DuckDuckGo)", "tavily": "网络搜索 (Tavily)"}


# ═══════════════════════════════════════════════════════════════════
# 公共 API
# ═══════════════════════════════════════════════════════════════════

def build_interactive_html(result: dict) -> str:
    """构建自包含交互式 HTML 页面，完整保存问答结果。

    包含: 答案 / 可视化图表 (搜索树或 DAG 轮次 Tab) / 检索来源 / 执行指标。
    所有 CSS/JS 内联，可在任意浏览器中离线打开。
    """
    question = result.get("answer", "")[:80] or "查询结果"
    algorithm = result.get("algorithm", "unknown")
    search_source = result.get("search_source")
    elapsed = result.get("elapsed", 0)
    answer = result.get("answer", "")
    chunks = result.get("chunks", []) or []

    algo_label = ALGO_LABELS.get(algorithm, algorithm)
    src_label = SOURCE_LABELS.get(search_source, str(search_source)) if search_source else "无"

    answer_html = _md_to_html(answer)
    graph_blobs: dict[str, str] = {}
    viz_html = _build_viz_section(result, graph_blobs)
    chunks_html = _build_chunks_html(chunks)

    search_count = None
    total_rounds = None
    term_reason = None
    pipeline_result = result.get("pipeline_result")
    if pipeline_result:
        search_count = pipeline_result.get("total_search_calls")
        total_rounds = pipeline_result.get("total_rounds")
        term_reason = pipeline_result.get("termination_reason")

    metrics_html = _build_metrics_bar(elapsed, len(chunks), search_count, total_rounds, term_reason)
    viz_js = _build_graph_js(graph_blobs)

    tab_js = ""
    if algorithm == "rag-loop" and pipeline_result:
        round_dags = pipeline_result.get("round_dags", [])
        if len(round_dags) > 1:
            tab_js = """function switchTab(idx) {
  var panels=document.querySelectorAll('.tab-panel');
  var btns=document.querySelectorAll('.tab-btn');
  panels.forEach(function(el,i){el.classList.toggle('active',i===idx);});
  btns.forEach(function(el,i){el.classList.toggle('active',i===idx);});
}"""

    return f"""<!DOCTYPE html>
<html lang="zh">
<head><meta charset="utf-8"><meta name="viewport" content="width=device-width,initial-scale=1">
<title>{_esc_html(question)}</title>
<style>
*{{margin:0;padding:0;box-sizing:border-box;}}
body{{font:14px/1.6 -apple-system,BlinkMacSystemFont,'Segoe UI','{FONT_CN}',sans-serif;background:#f5f7fa;color:#1f2937;padding:24px;max-width:1200px;margin:0 auto;}}
h1{{font-size:22px;color:#111827;margin-bottom:8px;}}
h2{{font-size:17px;color:#1f2937;margin:24px 0 12px;padding-bottom:6px;border-bottom:2px solid {ACCENT};}}
.badge{{display:inline-block;padding:2px 10px;border-radius:12px;font-size:12px;font-weight:600;margin-right:6px;}}
.badge-algo{{background:#dbeafe;color:#1e40af;}}
.badge-src{{background:#f3f4f6;color:#6b7280;}}
.meta-line{{color:#9ca3af;font-size:12px;margin:8px 0;}}
.answer{{background:#fff;padding:20px 24px;border-radius:10px;border:1px solid #e5e7eb;margin:12px 0;}}
.answer p{{margin:8px 0;}}
.answer code{{background:#f1f5f9;padding:2px 6px;border-radius:4px;font-size:13px;}}
.answer pre{{background:#1e293b;color:#e2e8f0;padding:12px 16px;border-radius:8px;overflow-x:auto;margin:8px 0;font-size:13px;}}
.answer pre code{{background:none;padding:0;color:inherit;}}
.answer strong{{color:#111827;}}
.answer em{{color:#374151;}}
.answer li{{margin:4px 0 4px 20px;}}
.answer h2{{font-size:18px;border:none;margin:16px 0 8px;padding:0;}}
.answer h3{{font-size:16px;border:none;margin:12px 0 6px;padding:0;}}
.answer h4{{font-size:14px;border:none;margin:10px 0 4px;padding:0;}}
.answer h5{{font-size:13px;border:none;margin:8px 0 4px;padding:0;}}
.answer blockquote{{border-left:3px solid #d1d5db;padding:4px 12px;margin:8px 0;color:#6b7280;}}
.answer hr{{border:none;border-top:1px solid #e5e7eb;margin:12px 0;}}
.graph-wrap{{background:#fff;border-radius:10px;border:1px solid #e5e7eb;padding:12px;margin:12px 0;overflow:hidden;}}
.graph-wrap iframe{{width:100%;border:none;border-radius:8px;}}
.chunk{{background:#fff;border:1px solid #e5e7eb;border-radius:8px;padding:14px 16px;margin:8px 0;}}
.chunk .ctitle{{font-weight:700;color:#1f2937;}}
.chunk .cmeta{{font-size:12px;color:#9ca3af;margin:4px 0;}}
.chunk .ctext{{font-size:13px;color:#4b5563;margin-top:6px;}}
.query-list{{list-style:decimal;padding-left:24px;}}
.query-list li{{margin:4px 0;font:13px monospace;color:#374151;background:#f9fafb;padding:4px 8px;border-radius:4px;}}
.metrics-bar{{display:flex;flex-wrap:wrap;gap:10px;margin:12px 0;}}
.metric-card{{background:#fff;border:1px solid #e5e7eb;border-radius:8px;padding:10px 16px;min-width:90px;}}
.metric-card .ml{{font-size:11px;color:#9ca3af;}}
.metric-card .mv{{font-size:18px;font-weight:700;color:{ACCENT};}}
.tabs{{margin:12px 0;}}
.tab-btns{{display:flex;gap:4px;flex-wrap:wrap;margin-bottom:8px;}}
.tab-btns button{{padding:6px 14px;border:1px solid #d1d5db;border-radius:6px;background:#fff;cursor:pointer;font-size:13px;color:#374151;transition:all .15s;}}
.tab-btns button:hover{{background:#f3f4f6;}}
.tab-btns button.active{{background:{ACCENT};color:#fff;border-color:{ACCENT};}}
.tab-panel{{display:none;}}
.tab-panel.active{{display:block;}}
.overview{{display:flex;flex-wrap:wrap;gap:10px;margin:8px 0 16px;}}
.overview .stat{{background:#f9fafb;border-radius:6px;padding:8px 14px;}}
.overview .stat .sn{{font-size:11px;color:#9ca3af;}}
.overview .stat .sv{{font-size:15px;font-weight:600;color:#1f2937;}}
</style></head>
<body>
<h1>{_esc_html(question)}</h1>
<div>
<span class="badge badge-algo">{_esc_html(algo_label)}</span>
<span class="badge badge-src">{_esc_html(src_label)}</span>
</div>
{metrics_html}

<h2>答案</h2>
<div class="answer">{answer_html}</div>

{viz_html}

<h2>检索来源 ({len(chunks)})</h2>
{chunks_html if chunks else '<p style="color:#9ca3af;">无检索来源</p>'}

<p class="meta-line">导出时间: {datetime.now().strftime('%Y-%m-%d %H:%M')} &middot; 算法: {_esc_html(algo_label)} &middot; 耗时: {elapsed:.1f}s</p>

<script>
{tab_js}
// ── 图表 Blob 嵌入 (base64 解码 → Blob URL → iframe src) ──
(function() {{
{viz_js}
}})();
</script>
</body></html>"""


def build_pptx_bytes(result: dict) -> bytes:
    """构建 PPTX 文件的 bytes。"""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    algorithm = result.get("algorithm", "unknown")
    search_source = result.get("search_source")
    question = _first_line(result.get("answer", "")[:80] or "查询")
    answer = result.get("answer", "")
    elapsed = result.get("elapsed", 0)
    chunks = result.get("chunks", []) or []
    rewritten = result.get("rewritten_queries") or []

    # 指标
    search_count = None
    pipeline_result = result.get("pipeline_result")
    if pipeline_result:
        search_count = pipeline_result.get("total_search_calls")

    # ── 1. 封面 ──
    _add_cover_slide(prs, question, algorithm, search_source)

    # ── 2. 答案 ──
    _add_answer_slide(prs, answer)

    # ── 3. 执行概览 ──
    _add_metrics_slide(prs, algorithm, search_source, elapsed, len(chunks), search_count)

    # ── 4. 检索来源 ──
    if chunks:
        _add_chunks_slide(prs, chunks)

    # ── 5. 可视化 (按算法差异) ──
    if algorithm in ("rag-with-judge", "rag-loop"):
        png_path = _capture_graph_png(result)
        _add_visualization_slide(prs, algorithm, png_path)
    elif algorithm == "naive-rag" and rewritten:
        _add_rewritten_queries_slide(prs, rewritten)

    # ── 6. 附录 ──
    if algorithm in ("rag-with-judge", "rag-loop"):
        _add_appendix_slide(prs)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


# ═══════════════════════════════════════════════════════════════════
# PPT 各页构建
# ═══════════════════════════════════════════════════════════════════

def _add_cover_slide(prs, question, algorithm, search_source):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # 背景色条
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(0.08))  # MSO_SHAPE.RECTANGLE = 1
    bar.fill.solid()
    bar.fill.fore_color.rgb = _hex_to_rgb(ACCENT)
    bar.line.fill.background()
    # 标题
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.3), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = question
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = _hex_to_rgb("#111827")
    _set_font(p, FONT_CN)
    # 副标题
    txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(4.8), Inches(10.3), Inches(1.2))
    tf2 = txBox2.text_frame
    algo_label = ALGO_LABELS.get(algorithm, algorithm)
    src_label = SOURCE_LABELS.get(search_source, str(search_source)) if search_source else "无检索"
    p2 = tf2.paragraphs[0]
    p2.text = f"{algo_label}  ·  {src_label}"
    p2.font.size = Pt(16)
    p2.font.color.rgb = _hex_to_rgb("#6b7280")
    _set_font(p2, FONT_CN)
    p3 = tf2.add_paragraph()
    p3.text = datetime.now().strftime("%Y-%m-%d %H:%M")
    p3.font.size = Pt(13)
    p3.font.color.rgb = _hex_to_rgb("#9ca3af")
    _set_font(p3, FONT_CN)


def _add_answer_slide(prs, answer):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_title(slide, "LLM 答案")
    # 去除 Markdown 标记
    clean = _strip_markdown(answer)
    txBox = slide.shapes.add_textbox(Inches(1.2), Inches(1.5), Inches(10.9), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = clean
    p.font.size = Pt(14)
    p.font.color.rgb = _hex_to_rgb("#1f2937")
    p.line_spacing = Pt(22)
    _set_font(p, FONT_CN)


def _add_metrics_slide(prs, algorithm, search_source, elapsed, chunk_count, search_count):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_title(slide, "执行概览")
    algo_label = ALGO_LABELS.get(algorithm, algorithm)
    src_label = SOURCE_LABELS.get(search_source, str(search_source)) if search_source else "无"
    metrics = [
        ("算法", algo_label),
        ("搜索源", src_label),
        ("耗时", f"{elapsed:.1f}s"),
        ("检索 Chunk 数", str(chunk_count)),
    ]
    if search_count is not None:
        metrics.append(("搜索次数", str(search_count)))
    # 2 列卡片布局
    card_w = Inches(5)
    card_h = Inches(1.6)
    start_x = Inches(1.2)
    start_y = Inches(1.6)
    gap_x = Inches(0.4)
    gap_y = Inches(0.3)
    for i, (label, value) in enumerate(metrics):
        col = i % 2
        row = i // 2
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)
        _add_metric_card(slide, x, y, card_w, card_h, label, value)


def _add_chunks_slide(prs, chunks):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_title(slide, "检索来源")
    top_n = chunks[:8]
    rows = len(top_n) + 1
    cols = 5
    table = slide.shapes.add_table(rows, cols, Inches(0.6), Inches(1.5), Inches(12.1), Inches(0.45 * rows)).table
    # 列宽
    widths = [Inches(0.4), Inches(2.8), Inches(3.0), Inches(0.7), Inches(5.2)]
    for i, w in enumerate(widths):
        table.columns[i].width = w
    # 表头
    headers = ["#", "标题", "来源文档", "分数", "内容摘要"]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = _hex_to_rgb("#6b7280")
            _set_font(p, FONT_CN)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _hex_to_rgb("#f3f4f6")
    # 数据行
    for i, chunk in enumerate(top_n):
        row_idx = i + 1
        vals = [
            str(i + 1),
            (chunk.get("chunk_title") or "?")[:60],
            (chunk.get("context_title") or "")[:50],
            f"{chunk.get('score', 0):.3f}" if chunk.get("score") else "N/A",
            (chunk.get("page_content") or "")[:200],
        ]
        for j, val in enumerate(vals):
            cell = table.cell(row_idx, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.color.rgb = _hex_to_rgb("#374151")
                _set_font(p, FONT_CN)


def _add_visualization_slide(prs, algorithm, png_path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if algorithm == "rag-with-judge":
        _slide_title(slide, "搜索过程可视化 — 搜索树")
    else:
        _slide_title(slide, "搜索过程可视化 — DAG 闭环图")

    if png_path and Path(png_path).exists():
        # 截图嵌入
        from PIL import Image
        img = Image.open(png_path)
        iw, ih = img.size
        # 缩放适配幻灯片 (最大约 10"×5")
        max_w, max_h = Inches(10.5), Inches(5.0)
        scale = min(max_w / iw, max_h / ih, 1.0)
        pw, ph = int(iw * scale), int(ih * scale)
        slide.shapes.add_picture(str(png_path), Inches(1.3), Inches(1.5), width=pw, height=ph)
        # 说明
        note_y = Inches(1.5) + ph + Inches(0.2)
    else:
        note_y = Inches(1.8)

    txBox = slide.shapes.add_textbox(Inches(1.3), note_y, Inches(10.4), Inches(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "💡 双击同目录下的 HTML 文件可在浏览器中查看可交互的搜索过程图（支持缩放/平移/点击节点）。"
    p.font.size = Pt(12)
    p.font.color.rgb = _hex_to_rgb("#9ca3af")
    _set_font(p, FONT_CN)


def _add_rewritten_queries_slide(prs, rewritten):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_title(slide, "多查询改写")
    txBox = slide.shapes.add_textbox(Inches(1.2), Inches(1.5), Inches(10.9), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, q in enumerate(rewritten):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{i + 1}.  {q}"
        p.font.size = Pt(14)
        p.font.color.rgb = _hex_to_rgb("#374151")
        p.line_spacing = Pt(24)
        _set_font(p, FONT_CN)


def _add_appendix_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_title(slide, "交互式图表说明")
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10.3), Inches(4.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    lines = [
        ("📂 文件说明", True, Pt(18)),
        ("", False, Pt(8)),
        ("导出时生成了两个文件：", False, Pt(14)),
        ("  · PPTX 文件 — 当前演示文稿，包含答案、指标和图表截图", False, Pt(14)),
        ("  · HTML 文件 — 自包含交互式页面，可在浏览器中打开", False, Pt(14)),
        ("", False, Pt(8)),
        ("🔍 如何使用 HTML 文件", True, Pt(18)),
        ("", False, Pt(8)),
        ("1. 在同目录下找到 .html 文件", False, Pt(14)),
        ("2. 双击用浏览器打开（Chrome / Edge 均可）", False, Pt(14)),
        ("3. 搜索树/DAG 图支持:", False, Pt(14)),
        ("    · 鼠标滚轮缩放", False, Pt(13)),
        ("    · 拖拽平移", False, Pt(13)),
        ("    · 点击节点查看详情", False, Pt(13)),
        ("", False, Pt(8)),
        ("📊 适合场景", True, Pt(18)),
        ("", False, Pt(8)),
        ("演示时: 用 PPTX 线性浏览", False, Pt(14)),
        ("深入讨论时: 打开 HTML 交互式探索搜索过程", False, Pt(14)),
    ]
    for text, bold, size in lines:
        if not tf.paragraphs[0].text:
            p = tf.paragraphs[0]
        elif text == "" and not bold:
            p = tf.add_paragraph()
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = size
        p.font.bold = bold
        if bold:
            p.font.color.rgb = _hex_to_rgb("#1f2937")
        else:
            p.font.color.rgb = _hex_to_rgb("#4b5563")
        _set_font(p, FONT_CN)


# ═══════════════════════════════════════════════════════════════════
# 图表截图 (Playwright)
# ═══════════════════════════════════════════════════════════════════

def _capture_graph_png(result: dict) -> str | None:
    """用 Playwright 将搜索树/DAG 渲染为 PNG，返回路径。失败返回 None。"""
    if not HAS_PLAYWRIGHT:
        return None
    try:
        graph_html = _render_graph_html(result)
        if not graph_html:
            return None
        tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
        tmp.close()
        with sync_playwright() as p:
            browser = p.chromium.launch(headless=True)
            page = browser.new_page(viewport={"width": 1400, "height": 900})
            page.set_content(graph_html, wait_until="networkidle")
            page.wait_for_timeout(500)
            page.screenshot(path=tmp.name, full_page=True)
            browser.close()
        return tmp.name
    except Exception:
        return None


def _render_graph_html(result: dict) -> str | None:
    """复用 graph_viewer 生成图表 HTML (完整独立页面, 供 Blob URL 嵌入)。"""
    from frontend.components.graph_viewer import _tree_to_graph, _build_iframe_html

    algorithm = result.get("algorithm", "")
    if algorithm == "rag-with-judge":
        search_path = result.get("search_path")
        if not search_path:
            return None
        nodes, edges = _tree_to_graph(search_path)
        title = search_path.get("question", "搜索树")[:60]
        graph_html, _ = _build_iframe_html(nodes, edges, "tree", title)
    elif algorithm == "rag-loop":
        pipeline = result.get("pipeline_result")
        if not pipeline:
            return None
        dag = pipeline.get("final_dag", {})
        nodes = dag.get("nodes", {})
        edges = dag.get("edges", [])
        if not nodes:
            return None
        graph_html, _ = _build_iframe_html(nodes, edges, "dag", "DAG 结构")
    else:
        return None

    # _build_iframe_html 已返回完整 HTML — 直接返回, 不做额外包裹
    return graph_html


# ═══════════════════════════════════════════════════════════════════
# HTML 构建辅助
# ═══════════════════════════════════════════════════════════════════

def _build_viz_section(result: dict, graph_blobs: dict) -> str:
    """构建可视化区域 HTML 并注册图表 Blob (base64 编码)。"""
    from frontend.components.graph_viewer import _tree_to_graph, _build_iframe_html

    algorithm = result.get("algorithm", "")
    rewritten = result.get("rewritten_queries") or []
    search_path = result.get("search_path")
    pipeline_result = result.get("pipeline_result")

    if algorithm == "naive-rag" and rewritten:
        items = "\n".join(f"<li>{_esc_html(q)}</li>" for q in rewritten)
        return f'<h2>多查询改写</h2><ol class="query-list">{items}</ol>'

    elif algorithm == "rag-with-judge":
        if not search_path:
            return '<p style="color:#9ca3af;">无搜索树数据</p>'
        nodes, edges = _tree_to_graph(search_path)
        if not nodes:
            return '<p style="color:#9ca3af;">搜索树为空</p>'
        graph_html, iframe_h = _build_iframe_html(nodes, edges, "tree",
            search_path.get("question", "搜索树")[:60])
        graph_blobs["graph-tree"] = base64.b64encode(
            graph_html.encode("utf-8")).decode("ascii")
        return f'''<h2>搜索过程可视化 &mdash; 搜索树</h2>
<div class="graph-wrap"><iframe id="graph-tree" style="height:{iframe_h}px;" title="搜索树"></iframe></div>'''

    elif algorithm == "rag-loop":
        if not pipeline_result:
            return '<p style="color:#9ca3af;">无 DAG 数据</p>'
        return _build_rag_loop_tabs(pipeline_result, graph_blobs)

    return ""


def _build_rag_loop_tabs(pipeline_result: dict, graph_blobs: dict) -> str:
    """构建 rag_loop 轮次 Tab 面板 + DAG 执行概览。"""
    from frontend.components.graph_viewer import _build_iframe_html

    final_dag = pipeline_result.get("final_dag", {})
    round_dags = pipeline_result.get("round_dags", [])
    total_rounds = pipeline_result.get("total_rounds", 0)
    total_searches = pipeline_result.get("total_search_calls", "?")
    term_reason = pipeline_result.get("termination_reason", "?")

    all_nodes = final_dag.get("nodes", {})
    node_count = len(all_nodes)
    solved = sum(1 for n in all_nodes.values() if n.get("status") == "solved")
    blocked = sum(1 for n in all_nodes.values() if n.get("status") == "blocked")

    term_labels = {
        "all_conditions_met": "全部条件满足",
        "max_rounds": "达到最大轮次",
        "planner_failure": "Planner 失败",
        "all_conditions_met_early": "提前完成",
    }
    term_label = term_labels.get(term_reason, term_reason)

    parts = ['<h2>DAG 闭环执行</h2>']
    parts.append(f'''<div class="overview">
<div class="stat"><div class="sn">总轮次</div><div class="sv">{total_rounds}</div></div>
<div class="stat"><div class="sn">搜索次数</div><div class="sv">{total_searches}</div></div>
<div class="stat"><div class="sn">节点数</div><div class="sv">{node_count}</div></div>
<div class="stat"><div class="sn">SOLVED</div><div class="sv" style="color:#16a34a">{solved}</div></div>
<div class="stat"><div class="sn">BLOCKED</div><div class="sv" style="color:#dc2626">{blocked}</div></div>
<div class="stat"><div class="sn">终止原因</div><div class="sv" style="font-size:13px">{_esc_html(term_label)}</div></div>
</div>''')

    dags_to_render = round_dags if round_dags else [final_dag]

    if len(dags_to_render) == 1:
        dag = dags_to_render[0]
        nodes = dag.get("nodes", {})
        edges = dag.get("edges", [])
        if nodes:
            graph_html, iframe_h = _build_iframe_html(nodes, edges, "dag", "DAG 结构")
            graph_blobs["graph-dag"] = base64.b64encode(
                graph_html.encode("utf-8")).decode("ascii")
            parts.append(f'<div class="graph-wrap"><iframe id="graph-dag" style="height:{iframe_h}px;" title="DAG 结构"></iframe></div>')
        else:
            parts.append('<p style="color:#9ca3af;">DAG 无节点</p>')
    else:
        tab_btns = []
        tab_panels = []
        for i, dag in enumerate(dags_to_render):
            round_num = dag.get("round_number", i + 1)
            nodes = dag.get("nodes", {})
            edges = dag.get("edges", [])
            active = " active" if i == 0 else ""
            tab_btns.append(
                f'<button class="tab-btn{active}" onclick="switchTab({i})">'
                f'Round {round_num}</button>'
            )
            if nodes:
                iframe_id = f"graph-dag-r{i}"
                graph_html, iframe_h = _build_iframe_html(
                    nodes, edges, "dag", f"DAG — Round {round_num}")
                graph_blobs[iframe_id] = base64.b64encode(
                    graph_html.encode("utf-8")).decode("ascii")
                panel_html = (
                    f'<div class="graph-wrap">'
                    f'<iframe id="{iframe_id}" style="height:{iframe_h}px;" '
                    f'title="Round {round_num}"></iframe></div>'
                )
            else:
                panel_html = '<p style="color:#9ca3af;">此轮无节点</p>'
            tab_panels.append(f'<div class="tab-panel{active}">{panel_html}</div>')

        parts.append('<div class="tabs">')
        parts.append('<div class="tab-btns">' + "\n".join(tab_btns) + '</div>')
        parts.append("\n".join(tab_panels))
        parts.append('</div>')

    return "\n".join(parts)


def _build_metrics_bar(elapsed: float, chunk_count: int,
                       search_count=None, total_rounds=None,
                       term_reason=None) -> str:
    """构建执行指标条 HTML。"""
    cards = [
        ("耗时", f"{elapsed:.1f}s"),
        ("检索 Chunks", str(chunk_count)),
    ]
    if search_count is not None:
        cards.append(("搜索次数", str(search_count)))
    if total_rounds is not None:
        cards.append(("总轮次", str(total_rounds)))

    items = "\n".join(
        f'<div class="metric-card"><div class="ml">{label}</div>'
        f'<div class="mv">{value}</div></div>'
        for label, value in cards
    )
    return f'<div class="metrics-bar">{items}</div>'


def _build_graph_js(graph_blobs: dict) -> str:
    """生成 base64 解码 → Blob URL → iframe src 的 JS 代码。

    atob() 返回 binary string (每字符 = 1 字节), 但 Blob(string) 按 UTF-16 解析,
    会破坏多字节 UTF-8 序列 (中文等)。必须先将 binary string 转成 Uint8Array。
    """
    parts = []
    for iframe_id, b64_html in graph_blobs.items():
        parts.append(
            f'(function(){{'
            f'var s=atob("{b64_html}");'
            f'var a=new Uint8Array(s.length);'
            f'for(var i=0;i<s.length;i++)a[i]=s.charCodeAt(i);'
            f'var b=new Blob([a],{{type:"text/html"}});'
            f'document.getElementById("{iframe_id}").src=URL.createObjectURL(b);'
            f'}})();'
        )
    return "\n".join(parts)


def _md_to_html(text: str) -> str:
    """将基本 Markdown 转为 HTML (纯 Python 实现, 无外部依赖)。

    先保护代码块/行内代码 → HTML 转义其余文本 → 恢复代码 → 应用 Markdown 模式。
    确保 LLM 输出中的任意 HTML 标签都被安全转义。
    """
    if not text:
        return ""
    # ── 1. 保护代码区域, 避免被 HTML 转义 ──
    code_blocks = []

    def _save_code_block(m):
        code_blocks.append(f'<pre><code>{_esc_html(m.group(2))}</code></pre>')
        return f"\x00CODEBLOCK{len(code_blocks) - 1}\x00"

    def _save_inline_code(m):
        code_blocks.append(f'<code>{_esc_html(m.group(1))}</code>')
        return f"\x00CODEBLOCK{len(code_blocks) - 1}\x00"

    text = re.sub(r'```(?:\w*)\n([\s\S]*?)```', _save_code_block, text)
    text = re.sub(r'`([^`]+)`', _save_inline_code, text)

    # 额外保护: 已经生成的 HTML 标签 (来自之前的处理)
    text = text.replace("<pre><code>", "\x00PRECODE\x00")
    text = text.replace("</code></pre>", "\x00ENDPRECODE\x00")

    # ── 2. HTML 转义 ──
    text = _esc_html(text)

    # 恢复保护区域
    text = text.replace("\x00PRECODE\x00", "<pre><code>")
    text = text.replace("\x00ENDPRECODE\x00", "</code></pre>")
    for i, val in enumerate(code_blocks):
        text = text.replace(f"\x00CODEBLOCK{i}\x00", val)

    # ── 3. Markdown 模式 ──
    # 粗体 + 斜体
    text = re.sub(r'\*\*\*(.+?)\*\*\*', r'<strong><em>\1</em></strong>', text)
    text = re.sub(r'\*\*(.+?)\*\*', r'<strong>\1</strong>', text)
    text = re.sub(r'\*(.+?)\*', r'<em>\1</em>', text)
    # 标题
    text = re.sub(r'^#### (.+)$', r'<h5>\1</h5>', text, flags=re.MULTILINE)
    text = re.sub(r'^### (.+)$', r'<h4>\1</h4>', text, flags=re.MULTILINE)
    text = re.sub(r'^## (.+)$', r'<h3>\1</h3>', text, flags=re.MULTILINE)
    text = re.sub(r'^# (.+)$', r'<h2>\1</h2>', text, flags=re.MULTILINE)
    # 引用 (注意: > 已被转义为 &gt;)
    text = re.sub(r'^&gt; (.+)$', r'<blockquote>\1</blockquote>', text, flags=re.MULTILINE)
    # 水平线
    text = re.sub(r'^[-*_]{3,}\s*$', r'<hr>', text, flags=re.MULTILINE)
    # 无序列表项
    text = re.sub(r'^[\s]*[-*+]\s+(.+)$', r'<li>\1</li>', text, flags=re.MULTILINE)
    # 有序列表项
    text = re.sub(r'^[\s]*\d+\.\s+(.+)$', r'<li>\1</li>', text, flags=re.MULTILINE)
    # 将连续的 <li> 包装为 <ul>
    text = re.sub(r'((?:<li>.*</li>\n?)+)', r'<ul>\1</ul>', text)
    # ── 4. 段落处理 ──
    paragraphs = text.split('\n\n')
    wrapped = []
    for p in paragraphs:
        p = p.strip()
        if not p:
            continue
        if p.startswith('<'):
            wrapped.append(p)
        else:
            p = p.replace('\n', '<br>')
            wrapped.append(f'<p>{p}</p>')
    return '\n'.join(wrapped)


def _build_chunks_html(chunks: list) -> str:
    if not chunks:
        return ""
    parts = []
    for i, c in enumerate(chunks[:15], 1):
        title = _esc_html(c.get("chunk_title") or "?")
        ctx = _esc_html(c.get("context_title") or "")
        score = c.get("score")
        score_str = f"{score:.3f}" if score else "N/A"
        content = _esc_html((c.get("page_content") or "")[:300])
        parts.append(
            f'<div class="chunk"><div class="ctitle">{i}. {title}</div>'
            f'<div class="cmeta">来源: {ctx} · 分数: {score_str}</div>'
            f'<div class="ctext">{content}</div></div>'
        )
    return "\n".join(parts)


# ═══════════════════════════════════════════════════════════════════
# PPT 辅助函数
# ═══════════════════════════════════════════════════════════════════

def _slide_title(slide, text):
    txBox = slide.shapes.add_textbox(Inches(1.2), Inches(0.5), Inches(11.0), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = _hex_to_rgb(ACCENT)
    _set_font(p, FONT_CN)
    # 分隔线
    line = slide.shapes.add_shape(1, Inches(1.2), Inches(1.15), Inches(11.0), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = _hex_to_rgb("#e5e7eb")
    line.line.fill.background()


def _add_metric_card(slide, x, y, w, h, label, value):
    """在幻灯片上添加一个指标卡片 (圆角矩形 + 文字)。"""
    import pptx.util
    # 背景卡片
    shape = slide.shapes.add_shape(
        5, x, y, w, h  # MSO_SHAPE.ROUNDED_RECTANGLE = 5
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb("#f9fafb")
    shape.line.color.rgb = _hex_to_rgb("#e5e7eb")
    shape.line.width = Pt(1)
    # 标签
    txBox = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.25), w - Inches(0.6), Inches(0.45))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(12)
    p.font.color.rgb = _hex_to_rgb("#6b7280")
    _set_font(p, FONT_CN)
    # 值
    txBox2 = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.7), w - Inches(0.6), Inches(0.6))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = value
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = _hex_to_rgb(ACCENT)
    _set_font(p2, FONT_CN)


def _hex_to_rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.lstrip("#")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def _set_font(paragraph, font_name):
    """设置段落中文字的字体（含东亚字体回退）。"""
    paragraph.font.name = font_name
    # 为每个 run 也设置
    for run in paragraph.runs:
        run.font.name = font_name


def _strip_markdown(text: str) -> str:
    """移除基本 Markdown 标记，保留纯文本。"""
    # 代码块
    text = re.sub(r'```[\s\S]*?```', '', text)
    # 行内代码
    text = re.sub(r'`([^`]+)`', r'\1', text)
    # 粗体/斜体
    text = re.sub(r'\*\*\*(.+?)\*\*\*', r'\1', text)
    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
    text = re.sub(r'\*(.+?)\*', r'\1', text)
    # 标题
    text = re.sub(r'^#{1,6}\s+', '', text, flags=re.MULTILINE)
    # 水平线
    text = re.sub(r'^[-*_]{3,}\s*$', '', text, flags=re.MULTILINE)
    # 无序列表标记
    text = re.sub(r'^[\s]*[-*+]\s+', '· ', text, flags=re.MULTILINE)
    # 有序列表标记
    text = re.sub(r'^[\s]*\d+\.\s+', '', text, flags=re.MULTILINE)
    # 清理多余空行
    text = re.sub(r'\n{3,}', '\n\n', text)
    return text.strip()


def _first_line(text: str) -> str:
    return text.split("\n")[0].strip()


def _esc_html(s: str) -> str:
    return s.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;").replace('"', "&quot;")